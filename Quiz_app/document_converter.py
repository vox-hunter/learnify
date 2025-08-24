# Document Converter Module
# Converts various document formats to PDF for optimal Gemini processing

import os
import io
import logging
from typing import Tuple, Optional
from PIL import Image
import tempfile

# Set up logging
logger = logging.getLogger(__name__)

def convert_to_pdf(file_content: bytes, filename: str) -> Tuple[Optional[bytes], Optional[str]]:
    """
    Convert various document formats to PDF for optimal Gemini processing.
    
    Args:
        file_content: The binary content of the file
        filename: Original filename for format detection
        
    Returns:
        tuple: (pdf_bytes, error_message)
            - pdf_bytes: PDF content as bytes or None if conversion failed
            - error_message: Error message or None if successful
    """
    if not filename:
        return None, "Filename is required for format detection"
    
    _, ext = os.path.splitext(filename.lower())
    
    try:
        # If it's already a PDF, return as-is
        if ext == '.pdf':
            return file_content, None
        
        # Convert based on file type
        if ext in ['.docx', '.doc']:
            return _convert_word_to_pdf(file_content, filename)
        elif ext in ['.pptx', '.ppt']:
            return _convert_powerpoint_to_pdf(file_content, filename)
        elif ext in ['.xlsx', '.xls']:
            return _convert_excel_to_pdf(file_content, filename)
        elif ext in ['.txt', '.md', '.markdown']:
            return _convert_text_to_pdf(file_content, filename)
        elif ext in ['.html', '.htm']:
            return _convert_html_to_pdf(file_content, filename)
        elif ext in ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.webp']:
            return _convert_image_to_pdf(file_content, filename)
        else:
            # For unsupported formats, return the original content
            # Gemini will handle it as binary
            return file_content, None
            
    except Exception as e:
        logger.error(f"Error converting {filename} to PDF: {e}")
        return None, f"Failed to convert {filename} to PDF: {str(e)}"

def _convert_word_to_pdf(file_content: bytes, filename: str) -> Tuple[Optional[bytes], Optional[str]]:
    """Convert Word document to PDF"""
    try:
        from docx import Document
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.pagesizes import letter
        from reportlab.lib.units import inch
        
        # Read Word document
        doc = Document(io.BytesIO(file_content))
        
        # Create PDF
        pdf_buffer = io.BytesIO()
        pdf_doc = SimpleDocTemplate(pdf_buffer, pagesize=letter)
        styles = getSampleStyleSheet()
        story = []
        
        # Add title
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=16,
            spaceAfter=30,
        )
        story.append(Paragraph(f"Document: {filename}", title_style))
        story.append(Spacer(1, 20))
        
        # Add paragraphs
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                p = Paragraph(paragraph.text, styles['Normal'])
                story.append(p)
                story.append(Spacer(1, 12))
        
        # Build PDF
        pdf_doc.build(story)
        pdf_bytes = pdf_buffer.getvalue()
        pdf_buffer.close()
        
        return pdf_bytes, None
        
    except ImportError:
        return None, "Required libraries not available for Word conversion (python-docx, reportlab)"
    except Exception as e:
        return None, f"Word conversion error: {str(e)}"

def _convert_powerpoint_to_pdf(file_content: bytes, filename: str) -> Tuple[Optional[bytes], Optional[str]]:
    """Convert PowerPoint presentation to PDF"""
    try:
        from pptx import Presentation
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.pagesizes import letter
        from reportlab.lib.units import inch
        
        # Read PowerPoint presentation
        prs = Presentation(io.BytesIO(file_content))
        
        # Create PDF
        pdf_buffer = io.BytesIO()
        pdf_doc = SimpleDocTemplate(pdf_buffer, pagesize=letter)
        styles = getSampleStyleSheet()
        story = []
        
        # Add title
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=18,
            spaceAfter=30,
        )
        story.append(Paragraph(f"Presentation: {filename}", title_style))
        story.append(Spacer(1, 20))
        
        # Add slides
        for slide_num, slide in enumerate(prs.slides, 1):
            # Slide header
            slide_header = ParagraphStyle(
                'SlideHeader',
                parent=styles['Heading2'],
                fontSize=14,
                spaceAfter=20,
            )
            story.append(Paragraph(f"Slide {slide_num}", slide_header))
            
            # Extract text from slide
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            
            if slide_text:
                for text in slide_text:
                    p = Paragraph(text, styles['Normal'])
                    story.append(p)
                    story.append(Spacer(1, 8))
            else:
                story.append(Paragraph("(No text content)", styles['Italic']))
            
            story.append(PageBreak())
        
        # Build PDF
        pdf_doc.build(story)
        pdf_bytes = pdf_buffer.getvalue()
        pdf_buffer.close()
        
        return pdf_bytes, None
        
    except ImportError:
        return None, "Required libraries not available for PowerPoint conversion (python-pptx, reportlab)"
    except Exception as e:
        return None, f"PowerPoint conversion error: {str(e)}"

def _convert_excel_to_pdf(file_content: bytes, filename: str) -> Tuple[Optional[bytes], Optional[str]]:
    """Convert Excel spreadsheet to PDF"""
    try:
        import openpyxl
        from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, PageBreak
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.lib.pagesizes import letter, landscape
        from reportlab.lib import colors
        
        # Read Excel workbook
        workbook = openpyxl.load_workbook(io.BytesIO(file_content))
        
        # Create PDF
        pdf_buffer = io.BytesIO()
        pdf_doc = SimpleDocTemplate(pdf_buffer, pagesize=landscape(letter))
        styles = getSampleStyleSheet()
        story = []
        
        # Add title
        story.append(Paragraph(f"Spreadsheet: {filename}", styles['Title']))
        story.append(Spacer(1, 20))
        
        # Process each worksheet
        for sheet_name in workbook.sheetnames:
            worksheet = workbook[sheet_name]
            
            # Sheet header
            story.append(Paragraph(f"Sheet: {sheet_name}", styles['Heading1']))
            story.append(Spacer(1, 12))
            
            # Get data from worksheet
            data = []
            max_rows = min(100, worksheet.max_row)  # Limit to 100 rows
            max_cols = min(20, worksheet.max_column)  # Limit to 20 columns
            
            for row in worksheet.iter_rows(min_row=1, max_row=max_rows, max_col=max_cols, values_only=True):
                row_data = [str(cell) if cell is not None else "" for cell in row]
                data.append(row_data)
            
            if data:
                # Create table
                table = Table(data)
                table.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                    ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('FONTSIZE', (0, 0), (-1, 0), 8),
                    ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                    ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                    ('GRID', (0, 0), (-1, -1), 1, colors.black),
                    ('FONTSIZE', (0, 1), (-1, -1), 6),
                ]))
                story.append(table)
            else:
                story.append(Paragraph("(Empty sheet)", styles['Italic']))
            
            story.append(PageBreak())
        
        # Build PDF
        pdf_doc.build(story)
        pdf_bytes = pdf_buffer.getvalue()
        pdf_buffer.close()
        
        return pdf_bytes, None
        
    except ImportError:
        return None, "Required libraries not available for Excel conversion (openpyxl, reportlab)"
    except Exception as e:
        return None, f"Excel conversion error: {str(e)}"

def _convert_text_to_pdf(file_content: bytes, filename: str) -> Tuple[Optional[bytes], Optional[str]]:
    """Convert text file to PDF"""
    try:
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.lib.pagesizes import letter
        
        # Decode text content
        try:
            text_content = file_content.decode('utf-8')
        except UnicodeDecodeError:
            try:
                text_content = file_content.decode('latin-1')
            except UnicodeDecodeError:
                text_content = file_content.decode('utf-8', errors='ignore')
        
        # Create PDF
        pdf_buffer = io.BytesIO()
        pdf_doc = SimpleDocTemplate(pdf_buffer, pagesize=letter)
        styles = getSampleStyleSheet()
        story = []
        
        # Add title
        story.append(Paragraph(f"Text File: {filename}", styles['Title']))
        story.append(Spacer(1, 20))
        
        # Add text content
        paragraphs = text_content.split('\n\n')
        for para in paragraphs:
            if para.strip():
                # Handle markdown-style headers
                if para.startswith('# '):
                    p = Paragraph(para[2:], styles['Heading1'])
                elif para.startswith('## '):
                    p = Paragraph(para[3:], styles['Heading2'])
                elif para.startswith('### '):
                    p = Paragraph(para[4:], styles['Heading3'])
                else:
                    p = Paragraph(para, styles['Normal'])
                story.append(p)
                story.append(Spacer(1, 12))
        
        # Build PDF
        pdf_doc.build(story)
        pdf_bytes = pdf_buffer.getvalue()
        pdf_buffer.close()
        
        return pdf_bytes, None
        
    except ImportError:
        return None, "Required libraries not available for text conversion (reportlab)"
    except Exception as e:
        return None, f"Text conversion error: {str(e)}"

def _convert_html_to_pdf(file_content: bytes, filename: str) -> Tuple[Optional[bytes], Optional[str]]:
    """Convert HTML file to PDF"""
    try:
        import weasyprint
        
        # Decode HTML content
        try:
            html_content = file_content.decode('utf-8')
        except UnicodeDecodeError:
            html_content = file_content.decode('utf-8', errors='ignore')
        
        # Convert to PDF using WeasyPrint
        pdf_buffer = io.BytesIO()
        weasyprint.HTML(string=html_content).write_pdf(pdf_buffer)
        pdf_bytes = pdf_buffer.getvalue()
        pdf_buffer.close()
        
        return pdf_bytes, None
        
    except ImportError:
        return None, "Required libraries not available for HTML conversion (weasyprint)"
    except Exception as e:
        return None, f"HTML conversion error: {str(e)}"

def _convert_image_to_pdf(file_content: bytes, filename: str) -> Tuple[Optional[bytes], Optional[str]]:
    """Convert image file to PDF"""
    try:
        from reportlab.platypus import SimpleDocTemplate, Image as RLImage, Spacer
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.lib.pagesizes import letter, A4
        from reportlab.lib.utils import ImageReader
        
        # Open image
        img_buffer = io.BytesIO(file_content)
        img = Image.open(img_buffer)
        
        # Create PDF
        pdf_buffer = io.BytesIO()
        pdf_doc = SimpleDocTemplate(pdf_buffer, pagesize=A4)
        styles = getSampleStyleSheet()
        story = []
        
        # Add title
        story.append(Paragraph(f"Image: {filename}", styles['Title']))
        story.append(Spacer(1, 20))
        
        # Calculate image size to fit page
        page_width, page_height = A4
        margin = 72  # 1 inch margins
        max_width = page_width - 2 * margin
        max_height = page_height - 200  # Leave space for title and margins
        
        # Scale image to fit
        img_width, img_height = img.size
        scale_w = max_width / img_width
        scale_h = max_height / img_height
        scale = min(scale_w, scale_h, 1.0)  # Don't scale up
        
        scaled_width = img_width * scale
        scaled_height = img_height * scale
        
        # Create image for PDF
        img_buffer.seek(0)  # Reset buffer position
        rl_img = RLImage(ImageReader(img_buffer), width=scaled_width, height=scaled_height)
        story.append(rl_img)
        
        # Build PDF
        pdf_doc.build(story)
        pdf_bytes = pdf_buffer.getvalue()
        pdf_buffer.close()
        
        return pdf_bytes, None
        
    except ImportError:
        return None, "Required libraries not available for image conversion (PIL, reportlab)"
    except Exception as e:
        return None, f"Image conversion error: {str(e)}"

def should_convert_to_pdf(filename: str) -> bool:
    """
    Determine if a file should be converted to PDF for better Gemini processing.
    """
    if not filename:
        return False
    
    _, ext = os.path.splitext(filename.lower())
    
    # Convert these formats to PDF for better Gemini processing
    convertible_formats = {
        '.docx', '.doc',        # Word documents
        '.pptx', '.ppt',        # PowerPoint presentations
        '.xlsx', '.xls',        # Excel spreadsheets
        '.txt', '.md',          # Text files
        '.html', '.htm',        # HTML files
        '.png', '.jpg', '.jpeg', # Images
        '.gif', '.bmp', '.tiff', '.webp'
    }
    
    return ext in convertible_formats

def get_conversion_info(filename: str) -> str:
    """
    Get information about how a file will be processed.
    """
    if not filename:
        return "Unknown file type"
    
    _, ext = os.path.splitext(filename.lower())
    
    if ext == '.pdf':
        return "✅ PDF - Native Gemini document understanding with vision capabilities"
    elif should_convert_to_pdf(filename):
        return "🔄 Will be converted to PDF for optimal Gemini document understanding"
    else:
        return "📦 Will be processed as binary data by Gemini AI"
