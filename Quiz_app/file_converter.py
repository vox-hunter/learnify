# File Converter - Convert various file types to PDF for Gemini AI processing
# This ensures optimal document understanding with Gemini's vision capabilities

import os
import tempfile
import logging
from typing import Tuple, Optional
import io

logger = logging.getLogger(__name__)

def convert_to_pdf(file_content: bytes, filename: str) -> Tuple[Optional[bytes], Optional[str]]:
    """
    Convert various file formats to PDF for optimal Gemini AI processing.
    
    Args:
        file_content: The binary content of the file
        filename: Original filename for format detection
        
    Returns:
        tuple: (pdf_bytes, error_message)
            - pdf_bytes: PDF content as bytes, or None if conversion failed
            - error_message: Error description, or None if successful
    """
    
    file_ext = os.path.splitext(filename.lower())[1]
    
    try:
        # Handle different file types
        if file_ext == '.pdf':
            # Already PDF, return as-is
            return file_content, None
            
        elif file_ext in ['.txt', '.md', '.markdown']:
            return _convert_text_to_pdf(file_content, filename)
            
        elif file_ext in ['.docx', '.doc']:
            return _convert_word_to_pdf(file_content, filename)
            
        elif file_ext in ['.pptx', '.ppt']:
            return _convert_powerpoint_to_pdf(file_content, filename)
            
        elif file_ext in ['.xlsx', '.xls']:
            return _convert_excel_to_pdf(file_content, filename)
            
        elif file_ext in ['.html', '.htm']:
            return _convert_html_to_pdf(file_content, filename)
            
        elif file_ext in ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.webp']:
            return _convert_image_to_pdf(file_content, filename)
            
        else:
            # For unsupported formats, create a simple PDF with file info
            return _create_info_pdf(filename, len(file_content))
            
    except Exception as e:
        logger.error(f"Error converting {filename} to PDF: {e}")
        return None, f"Conversion failed: {str(e)}"

def _convert_text_to_pdf(content: bytes, filename: str) -> Tuple[Optional[bytes], Optional[str]]:
    """Convert text files to PDF"""
    try:
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import letter
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.units import inch
        
        # Decode text content
        try:
            text = content.decode('utf-8')
        except UnicodeDecodeError:
            text = content.decode('latin-1', errors='ignore')
        
        buffer = io.BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=letter)
        styles = getSampleStyleSheet()
        
        # Create content
        story = []
        story.append(Paragraph(f"<b>File: {filename}</b>", styles['Title']))
        story.append(Spacer(1, 12))
        
        # Split text into paragraphs
        paragraphs = text.split('\n\n')
        for para in paragraphs:
            if para.strip():
                story.append(Paragraph(para.replace('\n', '<br/>'), styles['Normal']))
                story.append(Spacer(1, 6))
        
        doc.build(story)
        return buffer.getvalue(), None
        
    except ImportError:
        return None, "ReportLab library not available for text conversion"
    except Exception as e:
        return None, f"Text conversion error: {str(e)}"

def _convert_word_to_pdf(content: bytes, filename: str) -> Tuple[Optional[bytes], Optional[str]]:
    """Convert Word documents to PDF"""
    try:
        from docx import Document
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import letter
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.styles import getSampleStyleSheet
        
        # Load Word document
        doc_buffer = io.BytesIO(content)
        document = Document(doc_buffer)
        
        # Create PDF
        buffer = io.BytesIO()
        pdf_doc = SimpleDocTemplate(buffer, pagesize=letter)
        styles = getSampleStyleSheet()
        story = []
        
        story.append(Paragraph(f"<b>File: {filename}</b>", styles['Title']))
        story.append(Spacer(1, 12))
        
        # Extract paragraphs
        for para in document.paragraphs:
            if para.text.strip():
                story.append(Paragraph(para.text, styles['Normal']))
                story.append(Spacer(1, 6))
        
        pdf_doc.build(story)
        return buffer.getvalue(), None
        
    except ImportError:
        return None, "python-docx and ReportLab libraries not available for Word conversion"
    except Exception as e:
        return None, f"Word conversion error: {str(e)}"

def _convert_powerpoint_to_pdf(content: bytes, filename: str) -> Tuple[Optional[bytes], Optional[str]]:
    """Convert PowerPoint presentations to PDF"""
    try:
        from pptx import Presentation
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import letter
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
        from reportlab.lib.styles import getSampleStyleSheet
        
        # Load PowerPoint
        ppt_buffer = io.BytesIO(content)
        presentation = Presentation(ppt_buffer)
        
        # Create PDF
        buffer = io.BytesIO()
        pdf_doc = SimpleDocTemplate(buffer, pagesize=letter)
        styles = getSampleStyleSheet()
        story = []
        
        story.append(Paragraph(f"<b>Presentation: {filename}</b>", styles['Title']))
        story.append(Spacer(1, 12))
        
        # Extract slides
        for i, slide in enumerate(presentation.slides):
            story.append(Paragraph(f"<b>Slide {i+1}</b>", styles['Heading1']))
            story.append(Spacer(1, 6))
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    story.append(Paragraph(shape.text, styles['Normal']))
                    story.append(Spacer(1, 6))
            
            if i < len(presentation.slides) - 1:
                story.append(PageBreak())
        
        pdf_doc.build(story)
        return buffer.getvalue(), None
        
    except ImportError:
        return None, "python-pptx and ReportLab libraries not available for PowerPoint conversion"
    except Exception as e:
        return None, f"PowerPoint conversion error: {str(e)}"

def _convert_excel_to_pdf(content: bytes, filename: str) -> Tuple[Optional[bytes], Optional[str]]:
    """Convert Excel spreadsheets to PDF"""
    try:
        import pandas as pd
        from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.lib import colors
        from reportlab.lib.pagesizes import letter, landscape
        
        # Load Excel file
        excel_buffer = io.BytesIO(content)
        
        # Read all sheets
        excel_file = pd.ExcelFile(excel_buffer)
        
        # Create PDF
        buffer = io.BytesIO()
        pdf_doc = SimpleDocTemplate(buffer, pagesize=landscape(letter))
        styles = getSampleStyleSheet()
        story = []
        
        story.append(Paragraph(f"<b>Spreadsheet: {filename}</b>", styles['Title']))
        story.append(Spacer(1, 12))
        
        for sheet_name in excel_file.sheet_names:
            df = pd.read_excel(excel_file, sheet_name=sheet_name)
            
            story.append(Paragraph(f"<b>Sheet: {sheet_name}</b>", styles['Heading1']))
            story.append(Spacer(1, 6))
            
            # Convert DataFrame to table data
            if not df.empty:
                # Limit to first 20 rows and 10 columns to fit in PDF
                df_limited = df.head(20).iloc[:, :10]
                
                table_data = [df_limited.columns.tolist()]
                table_data.extend(df_limited.values.tolist())
                
                # Create table
                table = Table(table_data)
                table.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                    ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('FONTSIZE', (0, 0), (-1, 0), 8),
                    ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                    ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                    ('FONTSIZE', (0, 1), (-1, -1), 7),
                    ('GRID', (0, 0), (-1, -1), 1, colors.black)
                ]))
                
                story.append(table)
            else:
                story.append(Paragraph("(Empty sheet)", styles['Normal']))
                
            story.append(Spacer(1, 12))
        
        pdf_doc.build(story)
        return buffer.getvalue(), None
        
    except ImportError:
        return None, "pandas and ReportLab libraries not available for Excel conversion"
    except Exception as e:
        return None, f"Excel conversion error: {str(e)}"

def _convert_html_to_pdf(content: bytes, filename: str) -> Tuple[Optional[bytes], Optional[str]]:
    """Convert HTML files to PDF"""
    try:
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.lib.pagesizes import letter
        import re
        
        # Decode HTML content
        try:
            html = content.decode('utf-8')
        except UnicodeDecodeError:
            html = content.decode('latin-1', errors='ignore')
        
        # Simple HTML to text conversion
        # Remove HTML tags but keep basic structure
        text = re.sub(r'<script[^>]*>.*?</script>', '', html, flags=re.DOTALL)
        text = re.sub(r'<style[^>]*>.*?</style>', '', text, flags=re.DOTALL)
        text = re.sub(r'<[^>]+>', '', text)
        text = re.sub(r'\s+', ' ', text).strip()
        
        # Create PDF
        buffer = io.BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=letter)
        styles = getSampleStyleSheet()
        story = []
        
        story.append(Paragraph(f"<b>HTML File: {filename}</b>", styles['Title']))
        story.append(Spacer(1, 12))
        
        # Split into paragraphs
        paragraphs = text.split('\n')
        for para in paragraphs:
            if para.strip():
                story.append(Paragraph(para, styles['Normal']))
                story.append(Spacer(1, 6))
        
        doc.build(story)
        return buffer.getvalue(), None
        
    except ImportError:
        return None, "ReportLab library not available for HTML conversion"
    except Exception as e:
        return None, f"HTML conversion error: {str(e)}"

def _convert_image_to_pdf(content: bytes, filename: str) -> Tuple[Optional[bytes], Optional[str]]:
    """Convert images to PDF"""
    try:
        from PIL import Image
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import letter
        
        # Load image
        img_buffer = io.BytesIO(content)
        image = Image.open(img_buffer)
        
        # Create PDF
        buffer = io.BytesIO()
        c = canvas.Canvas(buffer, pagesize=letter)
        
        # Get page dimensions
        page_width, page_height = letter
        
        # Calculate image dimensions to fit page
        img_width, img_height = image.size
        aspect_ratio = img_width / img_height
        
        # Scale to fit page with margins
        max_width = page_width - 72  # 1 inch margins
        max_height = page_height - 144  # 2 inch margins (top and bottom)
        
        if aspect_ratio > 1:  # Wide image
            new_width = min(max_width, img_width)
            new_height = new_width / aspect_ratio
        else:  # Tall image
            new_height = min(max_height, img_height)
            new_width = new_height * aspect_ratio
        
        # Center the image
        x = (page_width - new_width) / 2
        y = (page_height - new_height) / 2
        
        # Save image to temporary buffer for PDF
        temp_img_buffer = io.BytesIO()
        image.save(temp_img_buffer, format='PNG')
        temp_img_buffer.seek(0)
        
        # Draw image on PDF
        c.drawImage(temp_img_buffer, x, y, width=new_width, height=new_height)
        
        # Add filename at bottom
        c.setFont("Helvetica", 10)
        c.drawString(72, 50, f"Image: {filename}")
        
        c.save()
        return buffer.getvalue(), None
        
    except ImportError:
        return None, "Pillow library not available for image conversion"
    except Exception as e:
        return None, f"Image conversion error: {str(e)}"

def _create_info_pdf(filename: str, file_size: int) -> Tuple[Optional[bytes], Optional[str]]:
    """Create a simple PDF with file information for unsupported formats"""
    try:
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import letter
        
        buffer = io.BytesIO()
        c = canvas.Canvas(buffer, pagesize=letter)
        
        # Add file information
        c.setFont("Helvetica-Bold", 16)
        c.drawString(72, 750, "File Information")
        
        c.setFont("Helvetica", 12)
        c.drawString(72, 700, f"Filename: {filename}")
        c.drawString(72, 680, f"File size: {file_size:,} bytes ({file_size/1024/1024:.2f} MB)")
        c.drawString(72, 660, f"File extension: {os.path.splitext(filename)[1]}")
        
        c.drawString(72, 620, "Note: This file format could not be converted to PDF.")
        c.drawString(72, 600, "The AI will process the original file directly.")
        
        c.save()
        return buffer.getvalue(), None
        
    except ImportError:
        return None, "ReportLab library not available for PDF creation"
    except Exception as e:
        return None, f"PDF creation error: {str(e)}"

def get_conversion_requirements(filename: str) -> list:
    """Get list of required libraries for converting a specific file type"""
    file_ext = os.path.splitext(filename.lower())[1]
    
    requirements = ['reportlab']  # Base requirement for PDF generation
    
    if file_ext in ['.docx', '.doc']:
        requirements.append('python-docx')
    elif file_ext in ['.pptx', '.ppt']:
        requirements.append('python-pptx')
    elif file_ext in ['.xlsx', '.xls']:
        requirements.extend(['pandas', 'openpyxl'])
    elif file_ext in ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.webp']:
        requirements.append('pillow')
    
    return requirements
